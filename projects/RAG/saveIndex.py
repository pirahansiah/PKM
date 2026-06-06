import os
import ollama
import chromadb
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm
# Import specialized parsers
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# --- CONFIGURATION ---
PROJECT_PATH = "/Volumes/4tb/2026-05-26/_Farshid"
DB_PATH = "/Volumes/4tb/2026-05-26/_Farshid/my_code_index"
COLLECTION_NAME = "project_docs"
EMBED_MODEL = "nomic-embed-text"

# 1. Initialize Persistent ChromaDB
client = chromadb.PersistentClient(path=DB_PATH)
collection = client.get_or_create_collection(name=COLLECTION_NAME)

# 2. Advanced Text Splitter
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1200, # Increased for Qwen's better logic
    chunk_overlap=200,
    add_start_index=True
)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext in ['.txt', '.md', '.py', '.cpp', '.h', '.json', '.yaml', '.yml']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif ext == '.docx':
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext in ['.ppt', '.pptx']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return text

def index_files():
    all_files = []
    EXTENSIONS = ['.py', '.docx', '.cpp', '.h', '.md', '.txt', '.pdf', '.pptx', '.json', '.yaml', '.yml']
    
    for root, _, files in os.walk(PROJECT_PATH):
        if any(x in root for x in ['.git', 'node_modules', '__pycache__', 'my_code_index_qwen']):
            continue
        for file in files:
            if any(file.endswith(ext) for ext in EXTENSIONS):
                all_files.append(os.path.join(root, file))

    print(f"Found {len(all_files)} files. Starting indexing...")

    for file_path in tqdm(all_files, desc="Indexing"):
        content = extract_text(file_path)
        if not content or len(content.strip()) < 10:
            continue
            
        chunks = text_splitter.split_text(content)
        
        # Check if file is already indexed to save time
        # We check the first chunk ID; if it exists, we skip the file
        first_chunk_id = f"{file_path}_0"
        existing = collection.get(ids=[first_chunk_id])
        if existing and existing['ids']:
            continue 

        for i, chunk in enumerate(chunks):
            chunk_id = f"{file_path}_{i}"
            
            # Generate embedding
            try:
                res = ollama.embeddings(model=EMBED_MODEL, prompt=chunk)
                collection.add(
                    ids=[chunk_id],
                    embeddings=[res["embedding"]],
                    documents=[chunk],
                    metadatas=[{"source": file_path}]
                )
            except Exception as e:
                print(f"Ollama Error on {file_path}: {e}")

if __name__ == "__main__":
    # Ensure Ollama is running and model is pulled
    # os.system(f"ollama pull {EMBED_MODEL}")
    index_files()
    print(f"✨ Indexing complete! Data stored at: {DB_PATH}")