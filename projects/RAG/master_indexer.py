"""
master_indexer.py  –  Farshid Brain: Ultimate RAG Indexer & Query Engine
=========================================================================
Enhancements over v1:
  • Richer metadata  : git_branch, line_count, char_count, language, section,
                       last_modified, file_size_kb, content_hash (SHA-256)
  • True dedup       : content_hash prevents re-indexing identical files even
                       if their path changes (copy/rename resilient)
  • Resume safety    : per-file state tracked in a lightweight SQLite ledger
                       (no Chroma round-trips needed for skip checks)
  • Robust extraction: per-format error isolation, encoding fallback chain,
                       binary-file guard, min-content threshold
  • Project-aware    : auto-detects project + sub-module from folder structure,
                       stored as separate metadata fields for precise filtering
  • Smart chunking   : language-aware separators for code vs prose
  • Query engine     : project/language/type filters, re-ranking by recency,
                       streaming-style progressive output, conversation history
  • Observability    : per-run stats, failed-file log, progress bar with ETA
"""

import os
import sys
import uuid
import time
import json
import sqlite3
import hashlib
import logging
import argparse
import textwrap
from datetime import datetime
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional

import ollama
import chromadb
from tqdm import tqdm

# ─── Optional imports (degrade gracefully) ────────────────────────────────────
try:
    from pypdf import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter, Language
    HAS_LANGCHAIN = True
except ImportError:
    HAS_LANGCHAIN = False

try:
    import subprocess
    _git_test = subprocess.run(["git", "--version"], capture_output=True)
    HAS_GIT = _git_test.returncode == 0
except Exception:
    HAS_GIT = False

# ══════════════════════════════════════════════════════════════════════════════
#  CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════
PROJECTS_ROOT    = "/Volumes/4tb/"
DB_PATH          = "/Users/farshid/projects/RAG/ULTIMATE_RAG_INDEX"
LEDGER_PATH      = os.path.join(DB_PATH, "indexer_ledger.sqlite")
FAILED_LOG_PATH  = os.path.join(DB_PATH, "failed_files.log")
COLLECTION_NAME  = "farshid_brain"

EMBED_MODEL      = "nomic-embed-text"
FAST_MODEL       = "qwen2.5-coder:7b"   # snappy answers
DEEP_MODEL       = "qwen2.5-coder:14b"  # thorough answers

CHUNK_SIZE       = 1500
CHUNK_OVERLAP    = 200
BATCH_SIZE       = 50
MAX_WORKERS      = 4          # Tune for M3 / Ollama concurrency
MAX_FILE_MB      = 20         # Skip files larger than this
MIN_CONTENT_LEN  = 30         # Discard near-empty extractions

# ─── Language detection map ────────────────────────────────────────────────
EXT_TO_LANGUAGE = {
    ".py": "python", ".js": "javascript", ".ts": "typescript",
    ".cpp": "cpp", ".c": "c", ".h": "cpp", ".hpp": "cpp",
    ".java": "java", ".go": "go", ".rs": "rust",
    ".rb": "ruby", ".php": "php", ".swift": "swift",
    ".sh": "bash", ".bash": "bash",
    ".md": "markdown", ".txt": "text",
    ".json": "json", ".yaml": "yaml", ".yml": "yaml", ".xml": "xml",
    ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
    ".html": "html", ".css": "css",
}

# All indexed extensions
ALL_EXTENSIONS = set(EXT_TO_LANGUAGE.keys())

# ══════════════════════════════════════════════════════════════════════════════
#  LOGGING
# ══════════════════════════════════════════════════════════════════════════════
logging.basicConfig(
    level=logging.WARNING,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("indexer")


# ══════════════════════════════════════════════════════════════════════════════
#  LEDGER  (SQLite – fast, zero-dependency, resume-safe)
# ══════════════════════════════════════════════════════════════════════════════
class Ledger:
    """Tracks indexed files by content hash. Survives crashes and reruns."""

    def __init__(self, path: str):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.con = sqlite3.connect(path, check_same_thread=False)
        self.con.execute("""
            CREATE TABLE IF NOT EXISTS indexed_files (
                content_hash  TEXT PRIMARY KEY,
                file_path     TEXT NOT NULL,
                chunk_count   INTEGER,
                indexed_at    TEXT
            )
        """)
        self.con.execute("""
            CREATE TABLE IF NOT EXISTS failed_files (
                file_path   TEXT PRIMARY KEY,
                reason      TEXT,
                failed_at   TEXT
            )
        """)
        self.con.commit()

    def is_indexed(self, content_hash: str) -> bool:
        row = self.con.execute(
            "SELECT 1 FROM indexed_files WHERE content_hash=?", (content_hash,)
        ).fetchone()
        return row is not None

    def mark_indexed(self, content_hash: str, file_path: str, chunk_count: int):
        self.con.execute(
            "INSERT OR REPLACE INTO indexed_files VALUES (?,?,?,?)",
            (content_hash, file_path, chunk_count, datetime.utcnow().isoformat()),
        )
        self.con.commit()

    def mark_failed(self, file_path: str, reason: str):
        self.con.execute(
            "INSERT OR REPLACE INTO failed_files VALUES (?,?,?)",
            (file_path, reason, datetime.utcnow().isoformat()),
        )
        self.con.commit()

    def stats(self) -> dict:
        total   = self.con.execute("SELECT COUNT(*) FROM indexed_files").fetchone()[0]
        chunks  = self.con.execute("SELECT SUM(chunk_count) FROM indexed_files").fetchone()[0] or 0
        failed  = self.con.execute("SELECT COUNT(*) FROM failed_files").fetchone()[0]
        return {"indexed_files": total, "total_chunks": chunks, "failed_files": failed}

    def failed_list(self) -> list[tuple]:
        return self.con.execute(
            "SELECT file_path, reason FROM failed_files ORDER BY failed_at DESC"
        ).fetchall()


# ══════════════════════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def sha256_file(path: str) -> Optional[str]:
    """SHA-256 of file bytes. Returns None on read error."""
    try:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for block in iter(lambda: f.read(65536), b""):
                h.update(block)
        return h.hexdigest()
    except Exception:
        return None


def git_branch(path: str) -> str:
    """Return current git branch for the repo containing `path`, or ''."""
    if not HAS_GIT:
        return ""
    try:
        result = subprocess.run(
            ["git", "-C", str(Path(path).parent), "rev-parse", "--abbrev-ref", "HEAD"],
            capture_output=True, text=True, timeout=3
        )
        return result.stdout.strip() if result.returncode == 0 else ""
    except Exception:
        return ""


def detect_project_and_module(file_path: str, root: str) -> tuple[str, str]:
    """
    Given /root/ProjectA/submodule/src/foo.py → ('ProjectA', 'submodule')
    Single-level files                         → ('ProjectA', '')
    """
    relative = os.path.relpath(file_path, root)
    parts = relative.split(os.sep)
    project = parts[0] if len(parts) >= 1 else "unknown"
    module  = parts[1] if len(parts) >= 3 else ""   # only set when there's depth
    return project, module


def get_splitter(language: str):
    """Return a language-aware splitter when langchain is available."""
    if not HAS_LANGCHAIN:
        # Fallback: naive splitter
        class NaiveSplitter:
            def split_text(self, text):
                size, overlap = CHUNK_SIZE, CHUNK_OVERLAP
                chunks, start = [], 0
                while start < len(text):
                    end = min(start + size, len(text))
                    chunks.append(text[start:end])
                    start += size - overlap
                return chunks
        return NaiveSplitter()

    lang_map = {
        "python": Language.PYTHON, "javascript": Language.JS,
        "typescript": Language.JS, "cpp": Language.CPP, "c": Language.C,
        "java": Language.JAVA, "go": Language.GO, "rust": Language.RUST,
        "ruby": Language.RUBY, "html": Language.HTML,
        "markdown": Language.MARKDOWN,
    }
    if language in lang_map:
        return RecursiveCharacterTextSplitter.from_language(
            language=lang_map[language],
            chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP
        )
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP
    )


# ══════════════════════════════════════════════════════════════════════════════
#  EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════

ENCODING_FALLBACKS = ["utf-8", "latin-1", "cp1252", "utf-16"]

def extract_text(file_path: str) -> Optional[str]:
    """Extract plain text from any supported file type. Returns None on failure."""
    ext = os.path.splitext(file_path)[1].lower()

    # ── Binary guard ──────────────────────────────────────────────────────────
    if os.path.getsize(file_path) > MAX_FILE_MB * 1024 * 1024:
        raise ValueError(f"File exceeds {MAX_FILE_MB} MB limit")

    # ── Code / text / data files ──────────────────────────────────────────────
    if ext in {".py",".cpp",".c",".h",".hpp",".js",".ts",".go",".rs",
               ".java",".php",".sh",".bash",".rb",".swift",
               ".txt",".md",".json",".yaml",".yml",".xml",".css",".html"}:
        for enc in ENCODING_FALLBACKS:
            try:
                with open(file_path, "r", encoding=enc, errors="strict") as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        # Last-resort: ignore bad bytes
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # ── PDF ────────────────────────────────────────────────────────────────────
    elif ext == ".pdf":
        if not HAS_PDF:
            raise ImportError("pypdf not installed")
        pages = []
        for page in PdfReader(file_path).pages:
            t = page.extract_text()
            if t:
                pages.append(t)
        return "\n".join(pages)

    # ── DOCX ───────────────────────────────────────────────────────────────────
    elif ext == ".docx":
        if not HAS_DOCX:
            raise ImportError("python-docx not installed")
        doc = DocxDocument(file_path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)

    # ── PPTX ───────────────────────────────────────────────────────────────────
    elif ext == ".pptx":
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed")
        lines = []
        for i, slide in enumerate(Presentation(file_path).slides):
            slide_lines = [f"[Slide {i+1}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            lines.extend(slide_lines)
        return "\n".join(lines)

    return None


# ══════════════════════════════════════════════════════════════════════════════
#  INDEXER
# ══════════════════════════════════════════════════════════════════════════════

class MasterIndexer:

    def __init__(self):
        os.makedirs(DB_PATH, exist_ok=True)
        self.chroma  = chromadb.PersistentClient(path=DB_PATH)
        self.col     = self.chroma.get_or_create_collection(
            name=COLLECTION_NAME,
            metadata={"hnsw:space": "cosine"}
        )
        self.ledger  = Ledger(LEDGER_PATH)
        self._lock   = __import__("threading").Lock()  # guards Chroma writes

    # ── Single-file pipeline ───────────────────────────────────────────────────

    def process_file(self, file_path: str) -> int:
        """Extract → hash → skip-if-seen → chunk → embed → store. Returns chunk count."""

        ext      = os.path.splitext(file_path)[1].lower()
        language = EXT_TO_LANGUAGE.get(ext, "text")

        # 1. Content hash (dedup across renames/copies)
        content_hash = sha256_file(file_path)
        if content_hash is None:
            self.ledger.mark_failed(file_path, "hash_failed")
            return 0

        if self.ledger.is_indexed(content_hash):
            return 0  # Already indexed – skip silently

        # 2. Extract
        try:
            content = extract_text(file_path)
        except Exception as e:
            self.ledger.mark_failed(file_path, str(e)[:200])
            return 0

        if not content or len(content.strip()) < MIN_CONTENT_LEN:
            self.ledger.mark_failed(file_path, "empty_or_too_short")
            return 0

        # 3. Rich metadata
        stat    = os.stat(file_path)
        project, module = detect_project_and_module(file_path, PROJECTS_ROOT)
        branch  = git_branch(file_path)
        meta_base = {
            "source":        file_path,
            "project":       project,
            "module":        module,
            "language":      language,
            "file_type":     ext,
            "git_branch":    branch,
            "last_modified": datetime.utcfromtimestamp(stat.st_mtime).isoformat(),
            "file_size_kb":  round(stat.st_size / 1024, 2),
            "line_count":    content.count("\n"),
            "char_count":    len(content),
            "content_hash":  content_hash,
        }

        # 4. Chunk (language-aware separators)
        splitter = get_splitter(language)
        chunks   = splitter.split_text(content)

        ids, embeds, metadatas, docs = [], [], [], []
        file_uuid = str(uuid.uuid5(uuid.NAMESPACE_URL, file_path))

        for i, chunk in enumerate(chunks):
            try:
                res = ollama.embeddings(model=EMBED_MODEL, prompt=chunk)
                ids.append(f"{file_uuid}_{i}")
                embeds.append(res["embedding"])
                docs.append(chunk)
                metadatas.append({**meta_base, "chunk_index": i, "total_chunks": len(chunks)})
            except Exception as e:
                log.warning("Embed failed chunk %d of %s: %s", i, file_path, e)
                continue

        # 5. Batch write (thread-safe)
        if ids:
            with self._lock:
                # Write in BATCH_SIZE sub-batches to avoid Chroma size limits
                for start in range(0, len(ids), BATCH_SIZE):
                    sl = slice(start, start + BATCH_SIZE)
                    self.col.add(
                        ids=ids[sl], embeddings=embeds[sl],
                        metadatas=metadatas[sl], documents=docs[sl]
                    )
            self.ledger.mark_indexed(content_hash, file_path, len(ids))

        return len(ids)

    # ── Full scan ──────────────────────────────────────────────────────────────

    def run(self, project_filter: Optional[str] = None):
        print(f"\n🔍  Scanning  {PROJECTS_ROOT}")
        if project_filter:
            print(f"    └─ Limiting to project: {project_filter}")

        files_to_process = []
        skip_dirs = {'.git', 'node_modules', '__pycache__', 'env', 'venv',
                     '.venv', 'ULTIMATE_RAG_INDEX', '.DS_Store', 'dist', 'build'}

        for root, dirs, files in os.walk(PROJECTS_ROOT):
            dirs[:] = [d for d in dirs if d not in skip_dirs]
            if project_filter:
                # Only descend into the target project folder
                rel = os.path.relpath(root, PROJECTS_ROOT)
                top = rel.split(os.sep)[0]
                if rel != "." and top != project_filter:
                    dirs[:] = []
                    continue
            for f in files:
                if any(f.lower().endswith(ext) for ext in ALL_EXTENSIONS):
                    full = os.path.join(root, f)
                    # Quick size guard before adding to queue
                    try:
                        if os.path.getsize(full) <= MAX_FILE_MB * 1024 * 1024:
                            files_to_process.append(full)
                    except OSError:
                        pass

        print(f"📂  {len(files_to_process)} eligible files found\n")

        total_chunks = 0
        t0 = time.time()

        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futures = {executor.submit(self.process_file, fp): fp for fp in files_to_process}
            with tqdm(total=len(futures), unit="file", ncols=90,
                      bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]") as pbar:
                for future in as_completed(futures):
                    try:
                        total_chunks += future.result()
                    except Exception as e:
                        fp = futures[future]
                        self.ledger.mark_failed(fp, f"executor: {e}")
                    pbar.update(1)

        elapsed = time.time() - t0
        stats   = self.ledger.stats()

        print(f"\n{'─'*55}")
        print(f"  ✅  New chunks indexed  : {total_chunks}")
        print(f"  📚  Total files indexed : {stats['indexed_files']}")
        print(f"  🧩  Total chunks stored : {stats['total_chunks']}")
        print(f"  ❌  Failed files        : {stats['failed_files']}")
        print(f"  ⏱   Elapsed             : {elapsed:.1f}s")
        print(f"{'─'*55}\n")

        if stats['failed_files']:
            print(f"  ⚠️  See failed files → {FAILED_LOG_PATH}")
            with open(FAILED_LOG_PATH, "w") as log_f:
                for fp, reason in self.ledger.failed_list():
                    log_f.write(f"{reason}\t{fp}\n")

    # ── Query engine ───────────────────────────────────────────────────────────

    def query(
        self,
        user_query: str,
        project_filter:  Optional[str] = None,
        module_filter:   Optional[str] = None,
        language_filter: Optional[str] = None,
        type_filter:     Optional[str] = None,
        n_results:       int = 20,
        deep:            bool = False,
        history:         Optional[list] = None,
    ) -> str:
        """
        Project-aware RAG query with optional metadata filters.

        Filters (all optional, combinable):
            project_filter   – e.g. "myWebsite"
            module_filter    – e.g. "cv_pipeline"
            language_filter  – e.g. "python"
            type_filter      – e.g. ".md"
        """
        # 1. Embed query
        try:
            q_emb = ollama.embeddings(model=EMBED_MODEL, prompt=user_query)["embedding"]
        except Exception as e:
            return f"[Embedding failed: {e}]"

        # 2. Build Chroma `where` clause (supports multiple filters)
        where_clauses = []
        if project_filter:
            where_clauses.append({"project":   {"$eq": project_filter}})
        if module_filter:
            where_clauses.append({"module":    {"$eq": module_filter}})
        if language_filter:
            where_clauses.append({"language":  {"$eq": language_filter}})
        if type_filter:
            where_clauses.append({"file_type": {"$eq": type_filter}})

        where = None
        if len(where_clauses) == 1:
            where = where_clauses[0]
        elif len(where_clauses) > 1:
            where = {"$and": where_clauses}

        # 3. Retrieve
        try:
            results = self.col.query(
                query_embeddings=[q_emb],
                n_results=min(n_results, self.col.count() or 1),
                where=where,
                include=["documents", "metadatas", "distances"]
            )
        except Exception as e:
            return f"[Chroma query failed: {e}]"

        docs_raw   = results["documents"][0]
        metas_raw  = results["metadatas"][0]
        dists_raw  = results["distances"][0]

        if not docs_raw:
            return "No relevant context found. Try broadening filters or re-indexing."

        # 4. Re-rank: blend cosine similarity + recency bonus
        def recency_score(meta: dict) -> float:
            try:
                mod = datetime.fromisoformat(meta.get("last_modified", "2000-01-01"))
                age_days = (datetime.utcnow() - mod).days
                return max(0.0, 1.0 - age_days / 365)  # 1.0 = today, 0.0 = 1yr+
            except Exception:
                return 0.0

        ranked = sorted(
            zip(docs_raw, metas_raw, dists_raw),
            key=lambda x: x[2] - 0.05 * recency_score(x[1])   # lower distance = better
        )

        # 5. Build context block with source headers
        context_parts = []
        seen_hashes = set()
        for doc, meta, dist in ranked:
            # Deduplicate chunks with identical content
            doc_hash = hashlib.md5(doc.encode()).hexdigest()
            if doc_hash in seen_hashes:
                continue
            seen_hashes.add(doc_hash)

            project  = meta.get("project", "?")
            module   = meta.get("module", "")
            lang     = meta.get("language", "")
            src      = meta.get("source", "?")
            modified = meta.get("last_modified", "?")[:10]
            relevance = round((1 - dist) * 100, 1)

            header = f"[{project}{'/'+module if module else ''} | {lang} | modified:{modified} | relevance:{relevance}%]\nFILE: {src}"
            context_parts.append(f"{header}\n{doc}")

        context = "\n\n".join(context_parts)

        # 6. Build project-aware system prompt
        project_label = project_filter or "all projects"
        system_prompt = textwrap.dedent(f"""
            You are an expert software engineer and technical writer working inside
            Farshid's personal knowledge base.

            Active project scope: {project_label}
            Retrieved {len(context_parts)} context chunks from the indexed SSD.

            Instructions:
            - Prioritise the provided FILE CONTEXT over your own training knowledge.
            - When referencing code, cite the exact source file path.
            - For CV/website tasks, preserve the original writing style found in context.
            - If the context is insufficient, say so explicitly rather than guessing.
            - Be concise but complete.
        """).strip()

        # 7. Optionally include conversation history
        prompt_parts = []
        if history:
            for turn in history[-4:]:   # last 4 turns for context window budget
                prompt_parts.append(f"Human: {turn['human']}\nAssistant: {turn['assistant']}")
        prompt_parts.append(f"CONTEXT FROM SSD:\n{context}\n\nUSER REQUEST: {user_query}")
        full_prompt = "\n\n".join(prompt_parts)

        model = DEEP_MODEL if deep else FAST_MODEL
        print(f"\n🤖  Using model: {model}  |  {len(context_parts)} chunks  |  project: {project_label}\n")

        # 8. Generate
        try:
            response = ollama.generate(
                model=model,
                system=system_prompt,
                prompt=full_prompt,
                options={"temperature": 0.2, "num_ctx": 8192},
            )
            return response["response"]
        except Exception as e:
            return f"[Generation failed: {e}]"


# ══════════════════════════════════════════════════════════════════════════════
#  CLI
# ══════════════════════════════════════════════════════════════════════════════

HELP_TEXT = """
Commands:
  /project <name>    Set project filter  (blank to clear)
  /module  <name>    Set module filter
  /lang    <name>    Set language filter (python, markdown, yaml, …)
  /type    <ext>     Set file-type filter (.py, .md, .pdf, …)
  /deep               Toggle deep model (14b) on/off
  /n       <number>  Set number of retrieved chunks (default 20)
  /filters           Show active filters
  /clear             Clear all filters + history
  /stats             Show ledger stats
  /failed            List failed files
  /help              Show this help
  exit / quit        Exit
"""


def cli_query_loop(indexer: MasterIndexer):
    print("\n" + "═"*55)
    print("  🧠  Farshid Brain  –  Query Mode")
    print("═"*55)
    print(HELP_TEXT)

    filters = dict(project_filter=None, module_filter=None,
                   language_filter=None, type_filter=None)
    n_results = 20
    deep      = False
    history   = []

    while True:
        try:
            raw = input("❓ Ask (or /help): ").strip()
        except (KeyboardInterrupt, EOFError):
            print("\nBye!")
            break

        if not raw:
            continue

        lower = raw.lower()

        if lower in ("exit", "quit"):
            print("Bye!")
            break

        # ── Commands ────────────────────────────────────────────────────────
        if lower == "/help":
            print(HELP_TEXT)
        elif lower == "/stats":
            s = indexer.ledger.stats()
            print(f"\n  Indexed files : {s['indexed_files']}")
            print(f"  Total chunks  : {s['total_chunks']}")
            print(f"  Failed files  : {s['failed_files']}\n")
        elif lower == "/failed":
            failed = indexer.ledger.failed_list()
            if not failed:
                print("  No failed files.\n")
            else:
                print(f"\n  {'REASON':<30} PATH")
                for fp, reason in failed[:20]:
                    print(f"  {reason[:30]:<30} {fp}")
                if len(failed) > 20:
                    print(f"  … and {len(failed)-20} more (see {FAILED_LOG_PATH})")
                print()
        elif lower.startswith("/project"):
            val = raw[8:].strip()
            filters["project_filter"] = val or None
            print(f"  ✔ project filter → {filters['project_filter'] or '(cleared)'}")
        elif lower.startswith("/module"):
            val = raw[7:].strip()
            filters["module_filter"] = val or None
            print(f"  ✔ module filter  → {filters['module_filter'] or '(cleared)'}")
        elif lower.startswith("/lang"):
            val = raw[5:].strip()
            filters["language_filter"] = val or None
            print(f"  ✔ language filter→ {filters['language_filter'] or '(cleared)'}")
        elif lower.startswith("/type"):
            val = raw[5:].strip()
            filters["type_filter"] = val or None
            print(f"  ✔ type filter    → {filters['type_filter'] or '(cleared)'}")
        elif lower.startswith("/n "):
            try:
                n_results = int(raw[3:].strip())
                print(f"  ✔ n_results → {n_results}")
            except ValueError:
                print("  Usage: /n <number>")
        elif lower == "/deep":
            deep = not deep
            print(f"  ✔ deep model → {'ON  (' + DEEP_MODEL + ')' if deep else 'OFF (' + FAST_MODEL + ')'}")
        elif lower == "/filters":
            print(f"\n  project  : {filters['project_filter'] or '—'}")
            print(f"  module   : {filters['module_filter']  or '—'}")
            print(f"  language : {filters['language_filter'] or '—'}")
            print(f"  type     : {filters['type_filter']    or '—'}")
            print(f"  n        : {n_results}")
            print(f"  deep     : {'yes  (' + DEEP_MODEL + ')' if deep else 'no   (' + FAST_MODEL + ')'}\n")
        elif lower == "/clear":
            filters = dict(project_filter=None, module_filter=None,
                           language_filter=None, type_filter=None)
            history.clear()
            n_results = 20
            deep = False
            print("  ✔ Filters + history cleared.")

        # ── Query ────────────────────────────────────────────────────────────
        else:
            answer = indexer.query(
                user_query=raw,
                **filters,
                n_results=n_results,
                deep=deep,
                history=history,
            )
            print(f"\n{answer}\n")
            history.append({"human": raw, "assistant": answer})


# ══════════════════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Farshid Brain – Ultimate RAG Indexer & Query Engine",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    sub = parser.add_subparsers(dest="cmd")

    # index sub-command
    idx_p = sub.add_parser("index", help="Scan and index files")
    idx_p.add_argument("--project", help="Only index a specific project folder")

    # query sub-command
    sub.add_parser("query", help="Interactive query loop")

    # stats sub-command
    sub.add_parser("stats", help="Show ledger stats and exit")

    args = parser.parse_args()

    indexer = MasterIndexer()

    if args.cmd == "index":
        print("💡 TIP: Run via 'caffeinate python master_indexer.py index' to keep Mac awake.")
        indexer.run(project_filter=getattr(args, "project", None))

    elif args.cmd == "stats":
        s = indexer.ledger.stats()
        print(f"\nIndexed files : {s['indexed_files']}")
        print(f"Total chunks  : {s['total_chunks']}")
        print(f"Failed files  : {s['failed_files']}\n")

    elif args.cmd == "query":
        cli_query_loop(indexer)

    else:
        # Legacy interactive mode (no sub-command)
        print("\n🧠  Farshid Brain  –  RAG Indexer v2\n")
        print("  python master_indexer.py index          # index everything")
        print("  python master_indexer.py index --project myWebsite")
        print("  python master_indexer.py query          # interactive query")
        print("  python master_indexer.py stats\n")


if __name__ == "__main__":
    main()